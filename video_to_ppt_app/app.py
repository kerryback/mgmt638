import os
import cv2
import numpy as np
from flask import Flask, render_template, request, send_file, jsonify
from pptx import Presentation
from pptx.util import Inches
from werkzeug.utils import secure_filename
import shutil

app = Flask(__name__)
app.config['MAX_CONTENT_LENGTH'] = 500 * 1024 * 1024  # 500MB max file size
app.config['UPLOAD_FOLDER'] = 'uploads'
app.config['OUTPUT_FOLDER'] = 'outputs'

# Create necessary directories
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)
os.makedirs(app.config['OUTPUT_FOLDER'], exist_ok=True)

def remove_logo(frame, logo_height_percent=8):
    """
    Remove logo from bottom by cropping the frame.

    Args:
        frame: Input frame (numpy array)
        logo_height_percent: Percentage of frame height to crop from bottom

    Returns:
        Cropped frame without logo area at bottom
    """
    height, width = frame.shape[:2]

    # Calculate crop height (remove only from bottom)
    crop_height = int(height * (1 - logo_height_percent / 100))

    # Crop the frame (keep full width, crop bottom)
    cropped = frame[:crop_height, :]

    return cropped

def extract_unique_frames(video_path, threshold=25, min_time_between_slides=2.0, remove_logo_flag=True):
    """
    Extract unique frames from video that represent distinct slides.

    Args:
        video_path: Path to video file
        threshold: Pixel difference threshold (0-255)
        min_time_between_slides: Minimum seconds between slide transitions
        remove_logo_flag: Whether to remove logo from bottom right corner

    Returns:
        List of frame images (numpy arrays)
    """
    cap = cv2.VideoCapture(video_path)
    if not cap.isOpened():
        raise ValueError(f"Could not open video file: {video_path}")

    fps = cap.get(cv2.CAP_PROP_FPS)
    min_frames_between_slides = int(fps * min_time_between_slides)

    frames = []
    prev_gray = None
    frame_count = 0
    last_saved_frame = -min_frames_between_slides

    while True:
        ret, frame = cap.read()
        if not ret:
            break

        # Remove logo if requested
        if remove_logo_flag:
            frame = remove_logo(frame)

        # Convert to grayscale for comparison
        gray = cv2.cvtColor(frame, cv2.COLOR_BGR2GRAY)

        # Always save the first frame
        if prev_gray is None:
            frames.append(frame.copy())
            prev_gray = gray
            last_saved_frame = frame_count
        else:
            # Check if enough time has passed since last slide
            if (frame_count - last_saved_frame) >= min_frames_between_slides:
                # Calculate difference from previous saved frame
                diff = cv2.absdiff(prev_gray, gray)
                mean_diff = np.mean(diff)

                # If significant change detected, save this frame
                if mean_diff > threshold:
                    frames.append(frame.copy())
                    prev_gray = gray
                    last_saved_frame = frame_count

        frame_count += 1

    cap.release()
    return frames

def create_powerpoint(frames, output_path):
    """
    Create PowerPoint presentation from list of frame images.

    Args:
        frames: List of frame images (numpy arrays)
        output_path: Path to save the PowerPoint file
    """
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Create temp directory for frame images
    temp_dir = os.path.join(app.config['OUTPUT_FOLDER'], 'temp_frames')
    os.makedirs(temp_dir, exist_ok=True)

    try:
        for idx, frame in enumerate(frames):
            # Save frame temporarily
            temp_image_path = os.path.join(temp_dir, f'frame_{idx}.jpg')
            cv2.imwrite(temp_image_path, frame)

            # Add blank slide
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # Get image dimensions
            height, width = frame.shape[:2]
            img_aspect = width / height

            # Calculate dimensions to fit slide while maintaining aspect ratio
            slide_width = prs.slide_width
            slide_height = prs.slide_height
            slide_aspect = slide_width / slide_height

            if img_aspect > slide_aspect:
                # Image is wider than slide
                pic_width = slide_width
                pic_height = int(slide_width / img_aspect)
            else:
                # Image is taller than slide
                pic_height = slide_height
                pic_width = int(slide_height * img_aspect)

            # Center the image
            left = (slide_width - pic_width) // 2
            top = (slide_height - pic_height) // 2

            # Add picture to slide
            slide.shapes.add_picture(temp_image_path, left, top,
                                    width=pic_width, height=pic_height)

        # Save presentation
        prs.save(output_path)

    finally:
        # Clean up temp directory
        if os.path.exists(temp_dir):
            shutil.rmtree(temp_dir)

@app.route('/')
def index():
    return render_template('index.html')

@app.route('/upload', methods=['POST'])
def upload():
    if 'video' not in request.files:
        return jsonify({'error': 'No file uploaded'}), 400

    file = request.files['video']
    if file.filename == '':
        return jsonify({'error': 'No file selected'}), 400

    if not file.filename.lower().endswith('.mp4'):
        return jsonify({'error': 'Only MP4 files are supported'}), 400

    try:
        # Save uploaded video
        filename = secure_filename(file.filename)
        video_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
        file.save(video_path)

        # Extract frames
        frames = extract_unique_frames(video_path)

        # Create PowerPoint
        output_filename = os.path.splitext(filename)[0] + '.pptx'
        output_path = os.path.join(app.config['OUTPUT_FOLDER'], output_filename)
        create_powerpoint(frames, output_path)

        # Clean up uploaded video
        os.remove(video_path)

        return jsonify({
            'success': True,
            'filename': output_filename,
            'slide_count': len(frames)
        })

    except Exception as e:
        return jsonify({'error': str(e)}), 500

@app.route('/download/<filename>')
def download(filename):
    filepath = os.path.join(app.config['OUTPUT_FOLDER'], filename)
    return send_file(filepath, as_attachment=True)

if __name__ == '__main__':
    app.run(debug=True, host='0.0.0.0', port=5000)
